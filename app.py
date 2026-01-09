import streamlit as st
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import io

# --- 辅助函数：判断颜色深浅 ---
def is_light_color(color_int):
    """判断颜色是否太浅（接近白色）"""
    if color_int is None: return False
    r = (color_int >> 16) & 0xFF
    g = (color_int >> 8) & 0xFF
    b = color_int & 0xFF
    # 计算亮度 (YIQ formula)
    brightness = (r * 299 + g * 587 + b * 114) / 1000
    return brightness > 200 # 阈值：大于200认为是亮色

# --- 辅助函数：设置字体和颜色 ---
def set_font_style(run, font_size, font_color_int, force_black_text):
    # 1. 字号
    run.font.size = Pt(font_size)
    
    # 2. 字体：强制微软雅黑
    run.font.name = "Microsoft YaHei"
    rPr = run.font._element.get_or_add_rPr()
    ea = rPr.get_or_add_ea()
    ea.set(qn('w:eastAsia'), 'Microsoft YaHei')
    
    # 3. 颜色处理 (关键修复：防止白字隐形)
    final_color = font_color_int
    
    if force_black_text:
        # 如果开启了强制黑字（用于白背景模式），且原文字是亮色，则强制变黑
        if is_light_color(font_color_int):
            run.font.color.rgb = RGBColor(0, 0, 0)
            return
            
    # 正常设置颜色
    try:
        r = (final_color >> 16) & 0xFF
        g = (final_color >> 8) & 0xFF
        b = final_color & 0xFF
        run.font.color.rgb = RGBColor(r, g, b)
    except:
        run.font.color.rgb = RGBColor(0, 0, 0) # 兜底黑色

# --- 核心转换逻辑 ---
def convert_pdf_to_ppt(uploaded_file, mode_strategy):
    uploaded_file.seek(0)
    doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
    prs = Presentation()
    
    # 尺寸初始化
    if len(doc) > 0:
        first_page = doc[0]
        prs.slide_width = Pt(first_page.rect.width)
        prs.slide_height = Pt(first_page.rect.height)

    progress_bar = st.progress(0)
    status_text = st.empty()
    total_pages = len(doc)

    # 策略解构
    # mode_strategy: 1=保留原背景图, 2=纯白背景+图文分离
    include_bg_image = (mode_strategy == 1)
    force_black_text = (mode_strategy == 2) # 如果是纯白背景，强制浅色文字变黑

    for i, page in enumerate(doc):
        progress_bar.progress((i + 1) / total_pages)
        status_text.text(f"正在重构第 {i+1} / {total_pages} 页 (图文分离)...")
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # --- A. 背景处理 ---
        if include_bg_image:
            # 策略1：整页截图做背景（最稳，但背景不可编辑）
            pix = page.get_pixmap(dpi=150)
            img_bytes = pix.tobytes("png")
            slide.shapes.add_picture(io.BytesIO(img_bytes), 0, 0, width=prs.slide_width, height=prs.slide_height)

        # --- B. 独立图片提取 (策略2时启用) ---
        # 如果我们不使用整页截图做背景，我们需要把PDF里的小插图一个个扣出来
        if not include_bg_image:
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    # 获取图片在页面上的位置
                    img_rects = page.get_image_rects(xref)
                    if not img_rects: continue

                    # 提取图片数据
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    for rect in img_rects:
                        if rect.width > 1 and rect.height > 1: # 忽略极小噪点
                            slide.shapes.add_picture(
                                io.BytesIO(image_bytes),
                                Pt(rect.x0), Pt(rect.y0),
                                width=Pt(rect.width), height=Pt(rect.height)
                            )
                except Exception:
                    pass # 图片提取失败跳过，保证程序不崩

        # --- C. 文字提取与重构 ---
        text_data = page.get_text("dict", flags=fitz.TEXT_PRESERVE_LIGATURES | fitz.TEXT_PRESERVE_WHITESPACE)
        
        for block in text_data["blocks"]:
            if block["type"] == 0:  # 文本
                for line in block["lines"]:
                    # 坐标与尺寸校验
                    x0, y0, x1, y1 = line["bbox"]
                    width = x1 - x0
                    height = y1 - y0
                    if width <= 0 or height <= 0: continue
                    
                    # 创建文本框
                    txBox = slide.shapes.add_textbox(Pt(x0), Pt(y0), Pt(width), Pt(height))
                    tf = txBox.text_frame
                    tf.word_wrap = False 
                    
                    p = tf.paragraphs[0]
                    
                    for span in line["spans"]:
                        text = span["text"]
                        # 过滤掉虽然存在但没内容的空字符
                        if not text.strip(): continue
                        
                        run = p.add_run()
                        run.text = text
                        
                        # 核心修正：字体 + 智能颜色
                        set_font_style(run, span["size"], span["color"], force_black_text)
                    
                    # 混合模式下的遮罩（可选，这里为了清爽先去掉）
                    # if include_bg_image:
                    #     txBox.fill.solid()
                    #     txBox.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    #     txBox.fill.transparency = 0.5

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# --- 页面 UI ---
st.set_page_config(page_title="PDF 转 PPT (V5.0 修复版)", layout="wide")
st.title("🛠️ PDF 转 PPT：文字修复版")

if 'ppt_data' not in st.session_state:
    st.session_state['ppt_data'] = None
if 'file_name' not in st.session_state:
    st.session_state['file_name'] = "fixed.pptx"

col1, col2 = st.columns([1, 2])

with col1:
    st.warning("模式选择 (解决全白问题)")
    mode = st.radio(
        "请选择转换策略：",
        (1, 2),
        format_func=lambda x: "策略 A: 保留原背景 (稳健)" if x == 1 else "策略 B: 纯白背景 + 智能黑字 (可编辑性最强)"
    )
    
    st.markdown("""
    **为什么之前是全白？**
    可能是因为原 PDF 是深色背景+浅色文字。如果去除背景，白色的字在白色 PPT 上就“隐身”了。
    
    **✅ 策略 B 的改进：**
    如果检测到文字是浅色的，会自动强制变成**黑色**，确保你能看见。同时会尝试提取独立插图。
    """)

with col2:
    uploaded_file = st.file_uploader("上传 PDF 文件", type=["pdf"])
    
    if uploaded_file:
        if st.button("🚀 开始修复并转换", type="primary"):
            try:
                with st.spinner("正在分析图层颜色并重构..."):
                    ppt_io = convert_pdf_to_ppt(uploaded_file, mode)
                    st.session_state['ppt_data'] = ppt_io
                    st.session_state['file_name'] = f"{uploaded_file.name.split('.')[0]}_v5.pptx"
                st.success("✅ 修复完成！颜色已智能校正。")
            except Exception as e:
                st.error(f"程序运行出错: {e}")
                st.write("建议检查 PDF 是否加密或为纯图片扫描件。")

    if st.session_state['ppt_data'] is not None:
        st.download_button(
            label="⬇️ 下载修复后的 PPT",
            data=st.session_state['ppt_data'],
            file_name=st.session_state['file_name'],
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
